"""
exports.py — Word, PDF y PowerPoint profesionales
"""
import io
import re
from datetime import datetime


def _limpiar(texto: str) -> str:
    texto = re.sub(r'\*\*(.+?)\*\*', r'\1', texto)
    texto = re.sub(r'\*(.+?)\*', r'\1', texto)
    texto = re.sub(r'#{1,6}\s+', '', texto)
    texto = re.sub(r'<[^>]+>', '', texto)
    texto = texto.replace('&', 'y')
    return texto.strip()


def _limpiar_rl(texto: str) -> str:
    texto = _limpiar(texto)
    texto = texto.replace('&', '&amp;').replace('"', '&quot;')
    texto = re.sub(r'<[^>]*>', '', texto)
    return texto.strip()


def _parsear_secciones(texto: str) -> list:
    """Retorna lista de (titulo, [lineas]) detectando secciones numeradas."""
    secciones = []
    titulo_actual = ""
    lineas_actuales = []
    for linea in texto.split('\n'):
        l = linea.strip()
        if re.match(r'^\d+\.\s+\S+', l):
            if titulo_actual:
                secciones.append((titulo_actual, lineas_actuales))
            titulo_actual = l
            lineas_actuales = []
        else:
            lineas_actuales.append(l)
    if titulo_actual:
        secciones.append((titulo_actual, lineas_actuales))
    return secciones


# ── WORD ──────────────────────────────────────────────────────────────────────

def generar_word(analisis: dict, resultados_cruce: list = None) -> bytes:
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()

    # Márgenes
    for section in doc.sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1.2)
        section.right_margin = Inches(1.2)

    AZUL = RGBColor(0x1a, 0x5f, 0xa8)
    GRIS = RGBColor(0x44, 0x44, 0x44)

    # Título principal
    p_titulo = doc.add_paragraph()
    run = p_titulo.add_run('⚖️  ANÁLISIS NORMATIVO')
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = AZUL

    # Línea separadora via borde inferior
    from docx.oxml import OxmlElement
    pPr = p_titulo._p.get_or_add_pPr()
    pBdr = OxmlElement('w:pBdr')
    bottom = OxmlElement('w:bottom')
    bottom.set(qn('w:val'), 'single')
    bottom.set(qn('w:sz'), '6')
    bottom.set(qn('w:color'), '1a5fa8')
    pBdr.append(bottom)
    pPr.append(pBdr)

    # Metadata
    doc.add_paragraph('')
    meta_items = [
        f"📄  Norma: {_limpiar(analisis.get('titulo', 'N/D'))}",
        f"🏢  Organismo: {analisis.get('organismo', 'N/D')}",
        f"⚡  Impacto: {analisis.get('impacto_principal', 'N/D')}",
        f"📅  Vigencia: {analisis.get('vigencia', 'N/D')}",
        f"🕐  Análisis: {datetime.now().strftime('%d/%m/%Y %H:%M')}",
    ]
    for item in meta_items:
        p = doc.add_paragraph()
        r = p.add_run(item)
        r.font.size = Pt(10)
        r.font.color.rgb = GRIS

    doc.add_paragraph('')

    # Secciones del análisis
    texto = _limpiar(analisis.get('analisis_completo', ''))
    secciones = _parsear_secciones(texto)

    if secciones:
        for titulo_s, lineas in secciones:
            # Título de sección
            h = doc.add_heading(titulo_s, level=1)
            for run in h.runs:
                run.font.color.rgb = AZUL
                run.font.size = Pt(13)

            for linea in lineas:
                if not linea:
                    continue
                if linea.startswith('- ') or linea.startswith('• '):
                    doc.add_paragraph(linea[2:], style='List Bullet')
                elif re.match(r'^[A-Z][A-Z\s]+:$', linea):
                    p = doc.add_paragraph()
                    r = p.add_run(linea)
                    r.bold = True
                    r.font.color.rgb = AZUL
                else:
                    doc.add_paragraph(linea)
    else:
        doc.add_paragraph(texto)

    # Cruce
    if resultados_cruce:
        doc.add_heading('RESULTADOS DEL CRUCE', level=1)
        enc = sum(1 for r in resultados_cruce if r["estado"] == "ENCUADRA")
        no  = sum(1 for r in resultados_cruce if r["estado"] == "NO ENCUADRA")
        aal = sum(1 for r in resultados_cruce if r["estado"] == "A ANALIZAR")
        doc.add_paragraph(f"Total: {len(resultados_cruce)} | ✅ Encuadran: {enc} | ❌ No: {no} | ⚠️ A analizar: {aal}")
        tabla = doc.add_table(rows=1, cols=4)
        tabla.style = 'Table Grid'
        for i, h in enumerate(['Artículo', 'NCM', 'Estado', 'Fundamento']):
            c = tabla.rows[0].cells[i]
            c.text = h
            c.paragraphs[0].runs[0].bold = True
        for r in resultados_cruce[:50]:
            fila = tabla.add_row()
            fila.cells[0].text = r.get('articulo', '')
            fila.cells[1].text = r.get('ncm', '')
            fila.cells[2].text = r.get('estado', '')
            fila.cells[3].text = r.get('fundamento', '')

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


# ── PDF ───────────────────────────────────────────────────────────────────────

def generar_pdf(analisis: dict, resultados_cruce: list = None) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.colors import HexColor, white
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable, Table, TableStyle
    from reportlab.lib.units import mm

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4,
                            leftMargin=22*mm, rightMargin=22*mm,
                            topMargin=20*mm, bottomMargin=20*mm)

    AZUL = HexColor('#1a5fa8')
    GRIS = HexColor('#444444')
    AZUL_CLAR = HexColor('#e8f4fd')

    styles = getSampleStyleSheet()
    s_titulo = ParagraphStyle('titulo', fontSize=20, textColor=AZUL, fontName='Helvetica-Bold', spaceAfter=4)
    s_meta   = ParagraphStyle('meta', fontSize=9, textColor=GRIS, spaceAfter=2, leading=13)
    s_h1     = ParagraphStyle('h1', fontSize=12, textColor=AZUL, fontName='Helvetica-Bold', spaceBefore=10, spaceAfter=4)
    s_body   = ParagraphStyle('body', fontSize=10, textColor=GRIS, spaceAfter=3, leading=15)
    s_bullet = ParagraphStyle('bullet', fontSize=10, textColor=GRIS, spaceAfter=3, leading=15, leftIndent=12)

    story = []
    story.append(Paragraph('ANALISIS NORMATIVO', s_titulo))
    story.append(HRFlowable(width='100%', thickness=2, color=AZUL, spaceAfter=6))

    story.append(Paragraph(f"Norma: {_limpiar_rl(analisis.get('titulo', 'N/D'))}", s_meta))
    story.append(Paragraph(f"Organismo: {analisis.get('organismo', 'N/D')} | Impacto: {analisis.get('impacto_principal', 'N/D')}", s_meta))
    story.append(Paragraph(f"Vigencia: {analisis.get('vigencia', 'N/D')} | Analizado: {datetime.now().strftime('%d/%m/%Y %H:%M')}", s_meta))
    story.append(Spacer(1, 8*mm))

    texto = _limpiar(analisis.get('analisis_completo', ''))
    secciones = _parsear_secciones(texto)

    if secciones:
        for titulo_s, lineas in secciones:
            story.append(Paragraph(_limpiar_rl(titulo_s), s_h1))
            story.append(HRFlowable(width='100%', thickness=0.5, color=AZUL, spaceAfter=4))
            for linea in lineas:
                if not linea:
                    story.append(Spacer(1, 2*mm))
                    continue
                linea_l = _limpiar_rl(linea)
                if not linea_l:
                    continue
                if linea.startswith('- ') or linea.startswith('• '):
                    story.append(Paragraph(f"• {_limpiar_rl(linea[2:])}", s_bullet))
                else:
                    story.append(Paragraph(linea_l, s_body))
    else:
        story.append(Paragraph(_limpiar_rl(texto[:5000]), s_body))

    doc.build(story)
    buf.seek(0)
    return buf.read()


# ── PPT ───────────────────────────────────────────────────────────────────────

def generar_ppt(analisis: dict, resultados_cruce: list = None) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    AZUL      = RGBColor(0x1a, 0x5f, 0xa8)
    AZUL_CLAR = RGBColor(0xd6, 0xe8, 0xf7)
    BLANCO    = RGBColor(0xff, 0xff, 0xff)
    GRIS_OSC  = RGBColor(0x22, 0x22, 0x22)
    GRIS_MED  = RGBColor(0x55, 0x55, 0x55)
    ACENTO    = RGBColor(0x00, 0xb0, 0xd8)

    blank = prs.slide_layouts[6]

    def shape_rect(slide, left, top, width, height, color):
        s = slide.shapes.add_shape(1,
            Inches(left), Inches(top), Inches(width), Inches(height))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
        return s

    def add_text(slide, texto, left, top, width, height,
                 size=14, bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
        tb = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = _limpiar(str(texto))[:300]
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        return tf

    def add_text_multiline(slide, lineas, left, top, width, height,
                           size=12, color=None, max_lineas=16):
        tb = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        primero = True
        count = 0
        for linea in lineas:
            linea = _limpiar(linea.strip())
            if not linea or count >= max_lineas:
                continue
            p = tf.paragraphs[0] if primero else tf.add_paragraph()
            primero = False
            run = p.add_run()
            if linea.startswith('- ') or linea.startswith('• '):
                run.text = f"  • {linea[2:120]}"
            else:
                run.text = linea[:120]
            run.font.size = Pt(size)
            if color:
                run.font.color.rgb = color
            count += 1

    # ── SLIDE 1: PORTADA ──────────────────────────────────────────────────────
    slide1 = prs.slides.add_slide(blank)
    w, h = prs.slide_width, prs.slide_height

    # Fondo completo azul
    shape_rect(slide1, 0, 0, 13.33, 7.5, AZUL)
    # Banda inferior acento
    shape_rect(slide1, 0, 6.8, 13.33, 0.7, ACENTO)
    # Rectángulo decorativo derecho
    shape_rect(slide1, 10.5, 0, 2.83, 7.5, RGBColor(0x12, 0x4a, 0x8a))

    add_text(slide1, '⚖️  ANÁLISIS NORMATIVO',
             0.6, 1.2, 9.5, 0.8, size=14, bold=False, color=ACENTO)

    titulo_norma = _limpiar(analisis.get('titulo', 'Análisis Normativo'))
    add_text(slide1, titulo_norma,
             0.6, 2.0, 9.5, 2.0, size=28, bold=True, color=BLANCO)

    add_text(slide1, f"Organismo: {analisis.get('organismo', '')}",
             0.6, 4.3, 9.0, 0.5, size=13, color=AZUL_CLAR)
    add_text(slide1, f"Impacto: {analisis.get('impacto_principal', '')}",
             0.6, 4.85, 9.0, 0.5, size=13, color=AZUL_CLAR)
    add_text(slide1, f"Vigencia: {analisis.get('vigencia', 'N/D')}",
             0.6, 5.4, 9.0, 0.5, size=13, color=AZUL_CLAR)
    add_text(slide1, datetime.now().strftime('%d/%m/%Y'),
             0.6, 6.85, 4.0, 0.4, size=11, color=BLANCO)

    # ── SLIDES DE SECCIONES ───────────────────────────────────────────────────
    texto = _limpiar(analisis.get('analisis_completo', ''))
    secciones = _parsear_secciones(texto)

    ICONOS = ['📋', '🔑', '⚙️', '⚠️', '📌', '✅', '❓']

    for idx, (titulo_s, lineas) in enumerate(secciones[:7]):
        slide = prs.slides.add_slide(blank)
        icono = ICONOS[idx] if idx < len(ICONOS) else '📄'

        # Fondo blanco
        shape_rect(slide, 0, 0, 13.33, 7.5, BLANCO)
        # Banda superior azul
        shape_rect(slide, 0, 0, 13.33, 1.1, AZUL)
        # Banda inferior acento
        shape_rect(slide, 0, 7.1, 13.33, 0.4, ACENTO)
        # Número de slide
        shape_rect(slide, 0, 0, 0.12, 7.5, ACENTO)

        # Título en banda azul
        add_text(slide, f"{icono}  {titulo_s}",
                 0.3, 0.15, 12.5, 0.8, size=18, bold=True, color=BLANCO)

        # Contenido
        lineas_filtradas = [l for l in lineas if l.strip()]
        add_text_multiline(slide, lineas_filtradas,
                           0.4, 1.3, 12.5, 5.6,
                           size=12, color=GRIS_OSC, max_lineas=18)

        # Norma en pie
        add_text(slide, _limpiar(analisis.get('titulo', ''))[:80],
                 0.3, 7.12, 10.0, 0.3, size=9, color=BLANCO)

    # ── SLIDE CRUCE (si existe) ────────────────────────────────────────────────
    if resultados_cruce:
        slide = prs.slides.add_slide(blank)
        shape_rect(slide, 0, 0, 13.33, 7.5, BLANCO)
        shape_rect(slide, 0, 0, 13.33, 1.1, AZUL)
        shape_rect(slide, 0, 7.1, 13.33, 0.4, ACENTO)
        shape_rect(slide, 0, 0, 0.12, 7.5, ACENTO)

        add_text(slide, '📊  RESULTADOS DEL CRUCE',
                 0.3, 0.15, 12.5, 0.8, size=18, bold=True, color=BLANCO)

        enc = sum(1 for r in resultados_cruce if r["estado"] == "ENCUADRA")
        no  = sum(1 for r in resultados_cruce if r["estado"] == "NO ENCUADRA")
        aal = sum(1 for r in resultados_cruce if r["estado"] == "A ANALIZAR")

        stats = [
            f"Total analizados: {len(resultados_cruce)}",
            f"🟢  Encuadran: {enc}",
            f"🔴  No encuadran: {no}",
            f"🟡  A analizar: {aal}",
        ]
        add_text_multiline(slide, stats, 0.4, 1.5, 6.0, 4.0,
                           size=20, color=GRIS_OSC, max_lineas=6)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
